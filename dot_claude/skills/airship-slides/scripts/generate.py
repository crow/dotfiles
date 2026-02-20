#!/usr/bin/env python3
"""
Generate branded Airship presentations using the master PPTX template.

Usage:
    python3 generate.py <input.json> <output.pptx>

Input JSON format:
{
    "slides": [
        {
            "layout": "cover",
            "title": "Presentation Title",
            "subtitle": "Subtitle text"
        },
        {
            "layout": "section",
            "title": "Chapter Title",
            "label": "SECTION 01"
        },
        {
            "layout": "feature_text",
            "title": "Big statement text here"
        },
        {
            "layout": "three_icon",
            "title": "Slide Title",
            "tag": "SET TAG IN CAPS",
            "columns": [
                {"icon": "path/to/icon.png", "subtitle": "Label", "body": "Description"},
                {"icon": "path/to/icon.png", "subtitle": "Label", "body": "Description"},
                {"icon": "path/to/icon.png", "subtitle": "Label", "body": "Description"}
            ]
        },
        {
            "layout": "four_icon",
            "title": "Slide Title",
            "tag": "SET TAG IN CAPS",
            "columns": [
                {"icon": "path/to/icon.png", "subtitle": "Label", "body": "Description"},
                {"icon": "path/to/icon.png", "subtitle": "Label", "body": "Description"},
                {"icon": "path/to/icon.png", "subtitle": "Label", "body": "Description"},
                {"icon": "path/to/icon.png", "subtitle": "Label", "body": "Description"}
            ]
        },
        {
            "layout": "three_image",
            "title": "Slide Title",
            "tag": "SET TAG IN CAPS",
            "columns": [
                {"image": "path/to/img.jpg", "subtitle": "Label", "body": "Description"},
                {"image": "path/to/img.jpg", "subtitle": "Label", "body": "Description"},
                {"image": "path/to/img.jpg", "subtitle": "Label", "body": "Description"}
            ]
        },
        {
            "layout": "four_image",
            "title": "Slide Title",
            "tag": "SET TAG IN CAPS",
            "columns": [
                {"image": "path/to/img.jpg", "subtitle": "Label", "body": "Description"},
                {"image": "path/to/img.jpg", "subtitle": "Label", "body": "Description"},
                {"image": "path/to/img.jpg", "subtitle": "Label", "body": "Description"},
                {"image": "path/to/img.jpg", "subtitle": "Label", "body": "Description"}
            ]
        },
        {
            "layout": "split_left",
            "title": "Slide Title",
            "tag": "SET TAG IN CAPS",
            "body": "Body text with bullets...",
            "image": "path/to/image.jpg"
        },
        {
            "layout": "split_right",
            "title": "Slide Title",
            "tag": "SET TAG IN CAPS",
            "body": "Body text...",
            "image": "path/to/image.jpg"
        },
        {
            "layout": "chart",
            "title": "Chart Title",
            "tag": "SET TAG IN CAPS",
            "body": "Description text",
            "callout": "Key takeaway highlight text"
        },
        {
            "layout": "card_grid",
            "title": "Slide Title",
            "tag": "SET TAG IN CAPS",
            "cards": [
                {"title": "Card Title", "body": "Card content"},
                {"title": "Card Title", "body": "Card content"},
                {"title": "Card Title", "body": "Card content"}
            ]
        },
        {
            "layout": "content",
            "title": "Slide Title",
            "tag": "SET TAG IN CAPS",
            "subtitle": "Optional subtitle",
            "body": "Body content text"
        },
        {
            "layout": "closing",
            "title": "Thank You"
        }
    ]
}
"""

import json
import sys
import os
from copy import deepcopy
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
except ImportError:
    print("Error: python-pptx is required. Install with: pip3 install python-pptx")
    sys.exit(1)

# Layout name mapping: friendly name -> PPTX layout index
LAYOUT_MAP = {
    # Cover slides (6 variants with different backgrounds)
    "cover": 0,          # CUSTOM_5 - title + subtitle + slide#
    "cover_2": 1,        # CUSTOM_5_2
    "cover_3": 2,        # CUSTOM_5_2_1
    "cover_4": 3,        # CUSTOM_5_1
    "cover_5": 4,        # CUSTOM_5_1_1
    "cover_6": 5,        # CUSTOM_5_1_1_1

    # Section dividers
    "feature_text": 6,   # SECTION_HEADER_4_1 - wide title only
    "feature_text_2": 7, # SECTION_HEADER_4_1_3
    "section": 8,        # SECTION_HEADER_4_1_2_2 - title + label + slide#
    "section_2": 9,      # SECTION_HEADER_4_1_2_1

    # Content slides
    "content": 10,       # ONE_COLUMN_TEXT - title + body + tag + subtitle
    "blank_title": 11,   # ONE_COLUMN_TEXT_7 - title + slide# (for tables/custom)
    "content_2": 12,     # ONE_COLUMN_TEXT_6
    "content_3": 13,     # ONE_COLUMN_TEXT_5
    "two_column": 14,    # ONE_COLUMN_TEXT_4_1_2 - two body columns
    "two_column_2": 15,  # ONE_COLUMN_TEXT_4_1_2_1
    "three_column": 16,  # ONE_COLUMN_TEXT_4_1_1_1 - three body columns
    "three_icon": 17,    # ONE_COLUMN_TEXT_4_1_1_1_2 - 3 icons + text

    # Grid slides
    "card_grid": 18,     # CUSTOM_8 - 3 card feature
    "three_image": 19,   # CUSTOM_8_2 - 3 images
    "four_image": 20,    # CUSTOM_8_2_1 - 4 images
    "cool_cards": 21,    # CUSTOM_8_1 - card feature with side content
    "four_icon": 22,     # ONE_COLUMN_TEXT_4_1_1_1_2_1 - 4 icons + text

    # Split slides
    "three_body": 23,    # ONE_COLUMN_TEXT_4_1_1_1_1
    "split_left": 24,    # CUSTOM_7 - text left, image right
    "split_right": 25,   # CUSTOM_7_1 - image left, text right
    "chart": 26,         # CUSTOM_7_1_1 - chart with callout

    # Closing
    "closing": 27,       # CUSTOM_6
}


def clear_slides(prs):
    """Remove all existing slides from the presentation, keeping layouts."""
    while len(prs.slides._sldIdLst) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[0]


def _placeholder_map(slide):
    """Build a dict mapping placeholder idx -> placeholder object."""
    return {ph.placeholder_format.idx: ph for ph in slide.placeholders}


def set_placeholder_text(slide, idx, text):
    """Safely set placeholder text, skipping if placeholder doesn't exist."""
    try:
        phs = _placeholder_map(slide)
        if idx in phs:
            phs[idx].text = text
            return True
    except Exception:
        pass
    return False


def set_placeholder_image(slide, idx, image_path):
    """Replace a picture placeholder with an image."""
    try:
        phs = _placeholder_map(slide)
        if idx in phs and os.path.exists(image_path):
            phs[idx].insert_picture(image_path)
            return True
    except Exception as e:
        print(f"  Warning: Could not set image at placeholder {idx}: {e}")
    return False


def add_slide(prs, slide_data):
    """Add a slide based on the layout type and content."""
    layout_name = slide_data.get("layout", "content")
    layout_idx = LAYOUT_MAP.get(layout_name)

    if layout_idx is None:
        print(f"  Warning: Unknown layout '{layout_name}', using 'content'")
        layout_idx = LAYOUT_MAP["content"]

    if layout_idx >= len(prs.slide_layouts):
        print(f"  Warning: Layout index {layout_idx} out of range, using 0")
        layout_idx = 0

    layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(layout)

    # Common fields
    title = slide_data.get("title", "")
    subtitle = slide_data.get("subtitle", "")
    tag = slide_data.get("tag", "")
    body = slide_data.get("body", "")
    callout = slide_data.get("callout", "")

    # ---- Cover slides (layouts 0-5) ----
    if layout_name.startswith("cover"):
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 1, subtitle)

    # ---- Section dividers (layouts 6-9) ----
    elif layout_name.startswith("section"):
        set_placeholder_text(slide, 0, title)
        label = slide_data.get("label", tag)
        set_placeholder_text(slide, 1, label)

    elif layout_name.startswith("feature_text"):
        set_placeholder_text(slide, 0, title)

    # ---- Content slides ----
    elif layout_name == "content" or layout_name.startswith("content_"):
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 1, body)
        set_placeholder_text(slide, 2, tag)
        if subtitle:
            set_placeholder_text(slide, 3, subtitle)

    elif layout_name == "blank_title":
        set_placeholder_text(slide, 0, title)

    # ---- Two/Three column text ----
    elif layout_name.startswith("two_column"):
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 2, tag)
        columns = slide_data.get("columns", [])
        if len(columns) >= 1:
            set_placeholder_text(slide, 1, columns[0].get("body", ""))
        if len(columns) >= 2:
            set_placeholder_text(slide, 3, columns[1].get("body", ""))

    elif layout_name == "three_column":
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 2, tag)
        columns = slide_data.get("columns", [])
        if len(columns) >= 1:
            set_placeholder_text(slide, 1, columns[0].get("body", ""))
        if len(columns) >= 2:
            set_placeholder_text(slide, 4, columns[1].get("body", ""))
        if len(columns) >= 3:
            set_placeholder_text(slide, 5, columns[2].get("body", ""))

    # ---- 3-icon grid (layout 17) ----
    elif layout_name == "three_icon":
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 2, tag)
        columns = slide_data.get("columns", [])
        # Placeholders: 3=pic, 4=pic, 5=pic, 6=subtitle, 8=subtitle, 13=subtitle
        # 1=body, 7=body, 9=body
        for i, col in enumerate(columns[:3]):
            icon = col.get("icon", "")
            sub = col.get("subtitle", "")
            bod = col.get("body", "")
            if i == 0:
                if icon: set_placeholder_image(slide, 3, icon)
                set_placeholder_text(slide, 6, sub)
                set_placeholder_text(slide, 1, bod)
            elif i == 1:
                if icon: set_placeholder_image(slide, 4, icon)
                set_placeholder_text(slide, 8, sub)
                set_placeholder_text(slide, 7, bod)
            elif i == 2:
                if icon: set_placeholder_image(slide, 5, icon)
                set_placeholder_text(slide, 13, sub)
                set_placeholder_text(slide, 9, bod)

    # ---- 4-icon grid (layout 22) ----
    elif layout_name == "four_icon":
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 2, tag)
        columns = slide_data.get("columns", [])
        # pic indices: 3, 6, 9, 15
        # subtitle indices: 4, 7, 13, 16
        # body indices: 1, 5, 8, 14
        pic_idx = [3, 6, 9, 15]
        sub_idx = [4, 7, 13, 16]
        bod_idx = [1, 5, 8, 14]
        for i, col in enumerate(columns[:4]):
            if col.get("icon"):
                set_placeholder_image(slide, pic_idx[i], col["icon"])
            set_placeholder_text(slide, sub_idx[i], col.get("subtitle", ""))
            set_placeholder_text(slide, bod_idx[i], col.get("body", ""))

    # ---- 3-image grid (layout 19) ----
    elif layout_name == "three_image":
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 1, tag)
        columns = slide_data.get("columns", [])
        # pic indices: 8, 9, 13
        # subtitle indices: 3, 5, 7
        # body indices: 2, 4, 6
        pic_idx = [8, 9, 13]
        sub_idx = [3, 5, 7]
        bod_idx = [2, 4, 6]
        for i, col in enumerate(columns[:3]):
            if col.get("image"):
                set_placeholder_image(slide, pic_idx[i], col["image"])
            set_placeholder_text(slide, sub_idx[i], col.get("subtitle", ""))
            set_placeholder_text(slide, bod_idx[i], col.get("body", ""))

    # ---- 4-image grid (layout 20) ----
    elif layout_name == "four_image":
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 1, tag)
        columns = slide_data.get("columns", [])
        pic_idx = [4, 7, 13, 16]
        sub_idx = [3, 6, 9, 15]
        bod_idx = [2, 5, 8, 14]
        for i, col in enumerate(columns[:4]):
            if col.get("image"):
                set_placeholder_image(slide, pic_idx[i], col["image"])
            set_placeholder_text(slide, sub_idx[i], col.get("subtitle", ""))
            set_placeholder_text(slide, bod_idx[i], col.get("body", ""))

    # ---- Card grid (layout 18) ----
    elif layout_name == "card_grid":
        set_placeholder_text(slide, 2, title)
        set_placeholder_text(slide, 1, tag)
        set_placeholder_text(slide, 3, subtitle)
        cards = slide_data.get("cards", [])
        # card titles: idx 0, 4, 5
        # card bodies: idx 6, 7, 8
        title_idx = [0, 4, 5]
        body_idx = [6, 7, 8]
        for i, card in enumerate(cards[:3]):
            set_placeholder_text(slide, title_idx[i], card.get("title", ""))
            set_placeholder_text(slide, body_idx[i], card.get("body", ""))

    # ---- Cool cards (layout 21) ----
    elif layout_name == "cool_cards":
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 14, tag)
        set_placeholder_text(slide, 13, body)
        cards = slide_data.get("cards", [])
        # card: pic=2,5,8  subtitle=3,6,9  body=1,4,7
        pic_idx = [2, 5, 8]
        sub_idx = [3, 6, 9]
        bod_idx = [1, 4, 7]
        for i, card in enumerate(cards[:3]):
            if card.get("icon"):
                set_placeholder_image(slide, pic_idx[i], card["icon"])
            set_placeholder_text(slide, sub_idx[i], card.get("subtitle", ""))
            set_placeholder_text(slide, bod_idx[i], card.get("body", ""))

    # ---- Split left (layout 24): text left, image right ----
    elif layout_name == "split_left":
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 1, body)
        set_placeholder_text(slide, 2, tag)
        image = slide_data.get("image", "")
        icon = slide_data.get("icon", "")
        if image: set_placeholder_image(slide, 4, image)
        if icon: set_placeholder_image(slide, 3, icon)

    # ---- Split right (layout 25): image left, text right ----
    elif layout_name == "split_right":
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 1, body)
        set_placeholder_text(slide, 3, tag)
        image = slide_data.get("image", "")
        icon = slide_data.get("icon", "")
        if image: set_placeholder_image(slide, 2, image)
        if icon: set_placeholder_image(slide, 4, icon)

    # ---- Chart with callout (layout 26) ----
    elif layout_name == "chart":
        set_placeholder_text(slide, 0, title)
        set_placeholder_text(slide, 1, body)
        set_placeholder_text(slide, 2, tag)
        set_placeholder_text(slide, 3, callout)

    # ---- Closing (layout 27) ----
    elif layout_name == "closing":
        set_placeholder_text(slide, 0, title)

    return slide


def generate(input_path, output_path, template_path=None):
    """Generate a presentation from JSON input."""
    # Find template
    if template_path is None:
        script_dir = Path(__file__).parent.parent
        template_path = script_dir / "assets" / "template" / "airship-master.pptx"

    if not os.path.exists(template_path):
        print(f"Error: Template not found at {template_path}")
        sys.exit(1)

    # Load input
    with open(input_path, 'r') as f:
        data = json.load(f)

    slides = data.get("slides", [])
    if not slides:
        print("Error: No slides in input JSON")
        sys.exit(1)

    # Load template and clear existing slides
    prs = Presentation(str(template_path))
    clear_slides(prs)

    print(f"Generating {len(slides)} slides...")

    for i, slide_data in enumerate(slides):
        layout_name = slide_data.get("layout", "content")
        print(f"  Slide {i+1}: {layout_name} - {slide_data.get('title', '')[:50]}")
        add_slide(prs, slide_data)

    prs.save(output_path)
    print(f"\n✅ Saved to {output_path}")
    print(f"   {len(slides)} slides generated using Airship brand template")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python3 generate.py <input.json> <output.pptx> [template.pptx]")
        print("\nSee script docstring for JSON input format.")
        sys.exit(1)

    input_path = sys.argv[1]
    output_path = sys.argv[2]
    template_path = sys.argv[3] if len(sys.argv) > 3 else None

    generate(input_path, output_path, template_path)
